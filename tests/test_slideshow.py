from unittest import TestCase
from datascience2pptx.slideshow import Slideshow
from datascience2pptx.slideshow import pptx_image
from pptx.dml.color import RGBColor
from matplotlib import pylab as P
import math
import pandas as pd
import numpy as np


class TestSlideshow(TestCase):
    """
    This tests are meant to work with TemplateExample.pptx
     (md5sum:  0fe375d702d14e2581e6132cc097d9d3)
    Which has the following layouts and placeholders:

        Portada
             Title 1
        Titulo
             Title 1
        Solo Imagen
             Picture Placeholder 1
        Titulo, Imagen
             Title 2
             Picture Placeholder 1
        Titulo, Tabla, dos Imagenes
             Title 4
             Picture Placeholder 1
             Picture Placeholder 2
             Table Placeholder 3
        Titulo, Cuatro Imagenes
             Title 3
             Picture Placeholder 1
             Picture Placeholder 2
             Picture Placeholder 4
             Picture Placeholder 5
        Titulo, dos imagenes
             Title 3
             Picture Placeholder 1
             Picture Placeholder 2
        Dos imagenes
             Picture Placeholder 1
             Picture Placeholder 2
        Imagen, Texto
             Picture Placeholder 1
             Content Placeholder 2
        Title, Content
             Title 1
             Text Placeholder 2
        Subtitulo
             Subtitle 1

    """

    def setup_slides(self):
        self.slide_image = {
            "name": "Slide1",
            "layout": "Solo Imagen",
            'placeholders': {
                "Picture Placeholder 1": self.the_I
            }
        }
        self.slide_dataframe = {
            "name": "Slide2",
            "layout": "Titulo, Tabla, dos Imagenes",
            'placeholders': {
                "Table Placeholder 3": self.the_df,
                "Title 4": self.the_string,
                "Picture Placeholder 1": self.the_I,
                "Picture Placeholder 2": self.the_fig
            }
        }
        self.slide_figure = {
            "name": "Slide3",
            "layout": "Solo Imagen",
            'placeholders': {
                "Picture Placeholder 1": self.the_fig
            }
        }
        self.slide_text = {
            "name": "Slide4",
            "layout": "Titulo",
            'placeholders': {
                "Title 1": self.the_string
            }
        }


    def setUp(self, templatepath="tests/TemplateExample.pptx",
              string2colormap={'▼': RGBColor(0, 255, 0),
                               '▲': RGBColor(255, 0, 0)}):
        self.sldsh = Slideshow(template=templatepath,
                               string2colormap=string2colormap)
        # for sl in self.sldsh.GetLayouts():
        #     print(sl.name)
        #     slide = self.sldsh.slideshow.slides.add_slide(sl)
        #     for ph in slide.placeholders:
        #         print("\t", ph.name)

        self.sldsh = Slideshow(template=templatepath,
                               string2colormap=string2colormap)


        # We initialize all the different objects to add
        self.the_string = "This a test string"

        X = list(range(30))
        Y = [math.cos(x) for x in X]
        fig = P.figure()
        P.plot(X, Y)
        self.the_fig = fig

        df = pd.DataFrame(np.random.randint(0, 10, size=(10, 4)), columns=list('ABCD'))
        df["textfield"] = "Some text"
        df["changes"] = ""
        i_row = df[df["A"] >= df["B"]].index
        i_row2 = df[df["A"] < df["B"]].index
        df.loc[i_row, 'changes'] = '▼'
        df.loc[i_row2, 'changes'] = '▲'
        self.the_df = df

        self.the_I = pptx_image("./tests/goose.jpg")

        self.setup_slides()

    def test__put_str(self):
        slides = [self.slide_text]

        self.sldsh.Execute(slides)
        print("put string")
        assert len(self.sldsh.slideshow.slides) == 1

    def test__put_figure(self):


        slides = [self.slide_figure]
        self.sldsh.Execute(slides)
        print("put figure")
        assert len(self.sldsh.slideshow.slides) == 1

    def test__put_dataframe(self):

        slides = [self.slide_dataframe]
        self.sldsh.Execute(slides)
        print("put dataframe")
        assert len(self.sldsh.slideshow.slides) > 0

    def test__put_bytes(self):
        assert True

    def test__put_image(self):
        slides = [self.slide_image]
        self.sldsh.Execute(slides)
        print("put figure")
        assert len(self.sldsh.slideshow.slides) == 1

    def test_SaveTo(self,path="/tmp/temp.pptx"):
        slides = [self.slide_text,
                  self.slide_figure,
                  self.slide_image,
                  self.slide_dataframe]
        self.sldsh.Execute(slides)
        self.sldsh.SaveTo(path)



    def tearDown(self):
        assert True


if __name__ == "__main__":
    test = TestSlideshow()

    test_methods = [method_name for method_name in dir(test)
                    if callable(getattr(test, method_name)) and
                    method_name.startswith("test_")]

    for m in test_methods:
        print("\nTesting", m)
        method = getattr(test, m)
        test.setUp()
        method()
        test.tearDown()
    test.tearDown()
